import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# AQUASENSE DEEP OCEANIC & ORBITAL EARTH OBSERVATION DESIGN SYSTEM
# ==============================================================================

# Background & Layer Hierarchy (Deep Abyssal to Oceanic Shelf)
OCEAN_ABYSS     = RGBColor(0x03, 0x07, 0x12) # Deep Space Abyssal Black (#030712)
OCEAN_DEEP      = RGBColor(0x07, 0x13, 0x26) # Mariana Oceanic Navy (#071326)
OCEAN_CARD      = RGBColor(0x0C, 0x1E, 0x3D) # Submerged Trench Card (#0C1E3D)
OCEAN_CARD_HI   = RGBColor(0x13, 0x2B, 0x54) # Continental Shelf Highlight (#132B54)
OCEAN_PANEL     = RGBColor(0x17, 0x35, 0x64) # Elevated Instrument Surface (#173564)

# Borders & Structural Grid
OCEAN_BORDER    = RGBColor(0x1D, 0x3D, 0x73) # Subtle Deep Sea Marine Border (#1D3D73)
OCEAN_BORDER_HI = RGBColor(0x02, 0x84, 0xC7) # Electric Azure Active Border (#0284C7)
BORDER_AQUA     = RGBColor(0x06, 0xD6, 0xA0) # Bioluminescent Emerald Teal (#06D6A0)

# Bioluminescent Accents & Spectral Palette
AQUA_BRIGHT     = RGBColor(0x22, 0xD3, 0xEE) # Phosphorescent Aqua Cyan (#22D3EE)
AQUA_ELECTRIC   = RGBColor(0x0E, 0xA5, 0xE9) # Electric Ocean Sky (#0EA5E9)
AQUA_DEEP       = RGBColor(0x02, 0x84, 0xC7) # Deep Atlantic Blue (#0284C7)
TEAL_BIO        = RGBColor(0x06, 0xD6, 0xA0) # Coastal Mangrove Bioluminescence (#06D6A0)
SOLAR_AMBER     = RGBColor(0xFB, 0xBF, 0x24) # Orbital Sunlight / Flare (#FBBF24)
CORAL_CRIMSON   = RGBColor(0xF4, 0x3F, 0x5E) # Desiccation / Loss Crimson (#F43F5E)
PURPLE_BIO      = RGBColor(0xA7, 0x8B, 0xFA) # Bio-Optical CDOM Purple (#A78BFA)

# Typography Hierarchy
TEXT_GLACIAL    = RGBColor(0xF0, 0xFD, 0xFA) # Pure Glacial White / Ice (#F0FDFA)
TEXT_FOAM       = RGBColor(0xCA, 0xDA, 0xEE) # Seafoam Silver Secondary (#CADDAE)
TEXT_SEABED     = RGBColor(0x73, 0x8C, 0xAD) # Seabed Muted Gray (#738CAD)
TEXT_MONO_CYAN  = RGBColor(0x38, 0xBD, 0xF8) # Telemetry Bright Cyan (#38BDF8)

FONT_TITLE = "Trebuchet MS"
FONT_BODY  = "Calibri"
FONT_MONO  = "Consolas"

# Assets Directory
ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "images")
IMG_SATELLITE   = os.path.join(ASSETS_DIR, "satellite_earth_orbit.jpg")
IMG_SPECTRAL    = os.path.join(ASSETS_DIR, "wetland_ndwi_spectral.jpg")
IMG_DRONE       = os.path.join(ASSETS_DIR, "multimodal_field_drone.jpg")
IMG_ARCH        = os.path.join(ASSETS_DIR, "technical_approach_diagram.jpg")
IMG_SCREENSHOT  = os.path.join(ASSETS_DIR, "current_app_screenshot.png")

IMG_NODE1       = os.path.join(ASSETS_DIR, "node1_stac_satellite.jpg")
IMG_NODE2       = os.path.join(ASSETS_DIR, "node2_band_math.jpg")
IMG_NODE3       = os.path.join(ASSETS_DIR, "node3_diff_mask.jpg")
IMG_NODE4       = os.path.join(ASSETS_DIR, "node4_gemini_ai.jpg")

# Initialize 16:9 Widescreen Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK_LAYOUT = prs.slide_layouts[6]

# ==============================================================================
# ROBUST DRAWING & LAYOUT HELPER FUNCTIONS
# ==============================================================================
def create_slide(prs, bg_color=OCEAN_ABYSS):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape

def add_pill_badge(slide, left, top, width, height, text, bg_color=OCEAN_CARD_HI, text_color=AQUA_BRIGHT, font_size=8.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = OCEAN_BORDER
    shape.line.width = Pt(0.75)
    
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = FONT_MONO
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_slide_header(slide, slide_num, category, title, subtitle=None):
    add_rect(slide, Inches(0.8), Inches(0.4), Inches(1.8), Pt(3), AQUA_BRIGHT)
    add_pill_badge(slide, Inches(0.8), Inches(0.52), Inches(2.2), Inches(0.28), f"ORBIT {slide_num:02d} // {category.upper()}", bg_color=OCEAN_CARD, text_color=SOLAR_AMBER, font_size=8.5)
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.82), Inches(11.7), Inches(0.65))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(23)
    p.font.bold = True
    p.font.color.rgb = TEXT_GLACIAL
    p.font.name = FONT_TITLE
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(11.0)
        p2.font.color.rgb = TEXT_SEABED
        p2.font.name = FONT_BODY
        p2.space_before = Pt(2)

def add_card(slide, left, top, width, height, title="", category_tag="", bg_color=OCEAN_CARD, border_color=OCEAN_BORDER, border_width=1):
    card = add_rect(slide, left, top, width, height, bg_color, border_color, border_width)
    
    if category_tag or title:
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.16), width - Inches(0.4), Inches(0.65))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        if category_tag:
            p0 = tf.paragraphs[0]
            p0.text = category_tag.upper()
            p0.font.size = Pt(8)
            p0.font.bold = True
            p0.font.color.rgb = AQUA_BRIGHT
            p0.font.name = FONT_MONO
            p0.space_after = Pt(2)
            
            if title:
                p1 = tf.add_paragraph()
                p1.text = title
                p1.font.size = Pt(11.5)
                p1.font.bold = True
                p1.font.color.rgb = TEXT_GLACIAL
                p1.font.name = FONT_TITLE
        elif title:
            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.size = Pt(11.5)
            p0.font.bold = True
            p0.font.color.rgb = TEXT_GLACIAL
            p0.font.name = FONT_TITLE
            
    return card

def add_bullet_list(slide, left, top, width, height, items, font_size=9.5, color=TEXT_FOAM, spacing=Pt(3.5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = spacing
    return txBox

def add_table_grid(slide, left, top, width, height, headers, rows, col_widths=None):
    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths and len(col_widths) == num_cols:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
            
    for c_idx, head in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = OCEAN_PANEL
        p = cell.text_frame.paragraphs[0]
        p.text = head.upper()
        p.font.size = Pt(8.0)
        p.font.bold = True
        p.font.color.rgb = AQUA_BRIGHT
        p.font.name = FONT_MONO
        p.alignment = PP_ALIGN.LEFT
        
    for r_idx, row_data in enumerate(rows):
        row_bg = OCEAN_CARD if r_idx % 2 == 0 else OCEAN_DEEP
        for c_idx, cell_data in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            p = cell.text_frame.paragraphs[0]
            p.text = str(cell_data)
            p.font.size = Pt(8.5)
            p.font.name = FONT_BODY
            p.font.color.rgb = TEXT_GLACIAL if c_idx == 0 else TEXT_FOAM
            if c_idx == 0:
                p.font.bold = True
                
    return table_shape

def add_image_with_frame(slide, image_path, left, top, width, height, label=None, border_color=OCEAN_BORDER_HI):
    if os.path.exists(image_path):
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        frame.fill.background()
        frame.line.color.rgb = border_color
        frame.line.width = Pt(1.5)
        
        if label:
            add_pill_badge(slide, left + Inches(0.12), top + Inches(0.12), Inches(len(label)*0.075 + 0.4), Inches(0.24), label, bg_color=OCEAN_ABYSS, text_color=AQUA_BRIGHT, font_size=7.0)
        return pic
    else:
        card = add_card(slide, left, top, width, height, title="Image Asset", category_tag=label or "GRAPHIC", border_color=border_color)
        return card

def add_footer_telemetry(slide, text="AQUASENSE OBSERVATORY • S2 L2A MSI + S1 C-SAR + DEM • 12-D SPECTRAL RAG • 10m/px • SOVEREIGN FEW-SHOT"):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.08), Inches(11.7), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7.5)
    p.font.color.rgb = TEXT_SEABED
    p.font.name = FONT_MONO
    p.alignment = PP_ALIGN.CENTER

# ==============================================================================
# SLIDE 1: COSMIC ORBITAL HERO SLIDE (WITH EMBEDDED SATELLITE IMAGE)
# ==============================================================================
slide1 = create_slide(prs)

add_pill_badge(slide1, Inches(0.8), Inches(0.6), Inches(3.8), Inches(0.28), "ORBITAL NODE: ESA S2 MSI + S1 C-SAR // STAC", bg_color=OCEAN_DEEP, text_color=AQUA_BRIGHT, font_size=8)
add_rect(slide1, Inches(0.8), Inches(0.95), Inches(1.8), Pt(4), AQUA_BRIGHT)

# Left Column: Observatory Title & Value Proposition
tx_title = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6.2), Inches(2.2))
tf = tx_title.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

p1 = tf.paragraphs[0]
p1.text = "AquaSense"
p1.font.size = Pt(50)
p1.font.bold = True
p1.font.color.rgb = TEXT_GLACIAL
p1.font.name = FONT_TITLE

p2 = tf.add_paragraph()
p2.text = "Multi-Sensor Planetary Hydro-Observatory & Sovereign RAG Engine"
p2.font.size = Pt(15)
p2.font.bold = True
p2.font.color.rgb = AQUA_BRIGHT
p2.font.name = FONT_BODY
p2.space_before = Pt(4)

p3 = tf.add_paragraph()
p3.text = "All-weather Sentinel-1 C-SAR radar fusion, 3D bathymetric volumetric retention (m³), bio-optical water quality (NDTI/NDCI), and sovereign 12-D spectral RAG intelligence."
p3.font.size = Pt(11.0)
p3.font.color.rgb = TEXT_FOAM
p3.font.name = FONT_BODY
p3.space_before = Pt(6)

# Core Feature Badges
badges_s1 = [
    ("🛰️ Dual Optical+SAR", OCEAN_PANEL, AQUA_BRIGHT),
    ("🌊 3D Bathymetry (m³)", OCEAN_PANEL, TEAL_BIO),
    ("💧 Bio-Optics (NDTI/NDCI)", OCEAN_PANEL, SOLAR_AMBER),
    ("🧬 Sovereign 12-D RAG", OCEAN_PANEL, PURPLE_BIO)
]
for idx, (txt, bg_c, txt_c) in enumerate(badges_s1):
    add_pill_badge(slide1, Inches(0.8 + idx * 1.55), Inches(3.7), Inches(1.48), Inches(0.32), txt, bg_color=bg_c, text_color=txt_c, font_size=7.5)

# Live Study Area Telemetry Panel (Left Lower)
add_card(slide1, Inches(0.8), Inches(4.3), Inches(6.0), Inches(2.5), title="Active Pilot: Pallikaranai Ramsar Basin", category_tag="MULTI-SENSOR LONGITUDINAL AUDIT", bg_color=OCEAN_CARD, border_color=OCEAN_BORDER_HI)
add_bullet_list(slide1, Inches(1.0), Inches(4.95), Inches(5.6), Inches(1.75), [
    "• Coordinates: [80.20°E, 12.91°N] to [80.23°E, 12.95°N] • Chennai, India",
    "• Multi-Sensor Constellation: Sentinel-2 MSI (10m) + Sentinel-1 RTC C-Band (5.405 GHz)",
    "• Volumetric Integration: Copernicus DEM GLO-30 Hypsometric Curves (V = ∫ A(z)dz)",
    "• Bio-Optical Water Quality: Turbidity (NDTI), Chlorophyll-a (NDCI), CDOM Carbon, WQI",
    "• Sovereign AI Engine: 12-D Feature Extractor + Few-Shot Classifier + Ramsar RAG Vector Store"
], font_size=9.0, spacing=Pt(2))

# Right Column: High-Resolution Satellite Visual Asset
add_image_with_frame(slide1, IMG_SATELLITE, Inches(7.1), Inches(0.8), Inches(5.4), Inches(4.0), label="SENTINEL-2 + SENTINEL-1 ORBITAL STREAM")

# Right Lower: Telemetry Specs Bar
add_card(slide1, Inches(7.1), Inches(5.0), Inches(5.4), Inches(1.8), title="Planetary Computer Multi-Sensor Status", category_tag="LIVE SENSOR TELEMETRY", bg_color=OCEAN_CARD, border_color=TEAL_BIO)
add_bullet_list(slide1, Inches(7.3), Inches(5.6), Inches(5.0), Inches(1.1), [
    "✔ STAC Nodes: Sentinel-2 L2A BOA & Sentinel-1 C-Band SAR RTC",
    "✔ Cloud Penetration: SAR Specular Backscatter (σ⁰_VV < -16 dB)",
    "✔ Local RAG Latency: <35ms deterministic in-memory vector lookup",
    "✔ Provenance Hash: 0x8a92f02c • Output: Immutable JSON Audit Bundle"
], font_size=8.5, spacing=Pt(2))

add_footer_telemetry(slide1)

# ==============================================================================
# SLIDE 2: THE PROBLEM STATEMENT — 4 CRITICAL FAILURES OF TRADITIONAL MONITORING
# ==============================================================================
slide2 = create_slide(prs)
add_slide_header(slide2, 2, "Problem Statement", "The Crisis: 4 Critical Failures in Traditional Water Monitoring", "Freshwater ecosystems are vanishing 3x faster than terrestrial forests, yet modern conservation is crippled by 4 fundamental technological blindspots.")

# 4 Clear & Distinct Problem Pillar Cards (2x2 Grid)
problems_grid = [
    ("1. THE OPTICAL CLOUD BLINDSPOT", "65%+ MONSOON OCCLUSION GAP", [
        "• Optical Satellites (Sentinel-2, Landsat) cannot penetrate dense monsoon clouds or night skies.",
        "• Tropical flood surges and peak runoff events occur exactly when optical observation is 100% blind.",
        "• Result: Emergency responders and hydrologists miss real-time flood crests and breach dynamics."
    ], CORAL_CRIMSON),
    ("2. THE 2D FLAT FOOTPRINT FALLACY", "IGNORING CUBIC STORAGE COLLAPSE", [
        "• Traditional GIS only calculates 2D surface area (km²), completely blind to basin sedimentation.",
        "• A lake can maintain 90% of its surface footprint while losing 60% of its actual cubic volume (m³).",
        "• Result: Municipalities experience sudden summer water crises despite 'healthy' 2D satellite maps."
    ], SOLAR_AMBER),
    ("3. INVISIBLE BIO-OPTICAL DEGRADATION", "TOXIC ALGAL & SILTATION DEFICIT", [
        "• Binary water index masks (NDWI > 0) classify pristine reservoirs and sewage lagoons identically.",
        "• Sediment plumes (Turbidity/TSS) and toxic cyanobacterial blooms (Chlorophyll-a) go undetected.",
        "• Result: Catastrophic fish kills and eutrophication occur without triggering any satellite warning."
    ], PURPLE_BIO),
    ("4. HIGH LATENCY & SOVEREIGN CLOUD RISK", "STATIC SURVEYS & API DEPENDENCY", [
        "• Manual field surveys are exorbitantly expensive (>$15k/basin) and obsolete upon delivery.",
        "• Decadal GIS updates (5-10 years) fail to detect weekly illegal infilling and wetland encroachment.",
        "• Heavy cloud LLMs introduce latency bottlenecks, API token costs, and sovereign data privacy risks."
    ], AQUA_DEEP)
]

for idx, (title, tag, pts, col) in enumerate(problems_grid):
    col_idx = idx % 2
    row_idx = idx // 2
    c_left = Inches(0.8 + col_idx * 5.95)
    c_top  = Inches(1.8 + row_idx * 2.5)
    
    add_card(slide2, c_left, c_top, Inches(5.75), Inches(2.35), title=title, category_tag=tag, bg_color=OCEAN_CARD, border_color=col, border_width=1.5)
    add_bullet_list(slide2, c_left + Inches(0.2), c_top + Inches(0.8), Inches(5.35), Inches(1.45), pts, font_size=9.0, spacing=Pt(3))

add_footer_telemetry(slide2)

# ==============================================================================
# SLIDE 3: OUR SOLUTION — 4 PILLARS OF AQUASENSE PLANETARY INTELLIGENCE
# ==============================================================================
slide3 = create_slide(prs)
add_slide_header(slide3, 3, "Our Solution", "The AquaSense Solution: 4 Pillars of Planetary Hydro-Intelligence", "An autonomous, multi-sensor observatory directly resolving all 4 traditional monitoring blindspots through orbital fusion and sovereign RAG AI.")

# 4 Clear & Distinct Solution Pillar Cards (2x2 Grid)
solutions_grid = [
    ("1. ALL-WEATHER SENTINEL-1 C-SAR FUSION", "100% CLOUD & NIGHT PENETRATION", [
        "• Emits 5.405 GHz (λ=5.6cm) active microwaves that penetrate 100% of dense clouds, fog, and night.",
        "• Water acts as a specular mirror reflecting pulses away (σ⁰_VV < -16 dB, crisp dark signature).",
        "• Dual Optical + SAR cross-validation ensures 365-day uninterrupted flood & desiccation tracking."
    ], TEAL_BIO),
    ("2. 3D HYPSOMETRIC BATHYMETRY (m³)", "VOLUMETRIC STORAGE & DEPTH STRATA", [
        "• Combines satellite surface footprint dynamics with Copernicus DEM GLO-30 / SRTM topography.",
        "• Hypsometric Integration: V(h) = ∫ A(z)dz calculates true cubic capacity (m³ & MCM).",
        "• Classifies depth strata: Littoral wetland (0-2m), submerged channel (2-5m), and deep core (>5m)."
    ], AQUA_BRIGHT),
    ("3. SPECTRAL BIO-OPTICAL WATER QUALITY", "TURBIDITY, CHL-A, CDOM & 0-100 WQI", [
        "• NDTI = (B04-B03)/(B04+B03): Quantifies sediment plumes and suspended solids (mg/L TSS, NTU).",
        "• NDCI = (B05-B04)/(B05+B04): Tracks algal bloom eutrophication via Carlson Trophic State (TSI).",
        "• Aggregated Water Quality Index (WQI 0-100) provides instantaneous ecological health scoring."
    ], SOLAR_AMBER),
    ("4. SOVEREIGN 12-D TENSOR & LOCAL RAG", "ZERO-CLOUD EDGE CLASSIFICATION (<35ms)", [
        "• 12-D Tensor: [B02..B08, NDWI, MNDWI, NDTI, NDCI, σ⁰_VV, σ⁰_VH, ΔDEM, WQI].",
        "• Few-Shot Logistic Regression & k-NN classifies land cover (water, wetland, built-up) in <35ms.",
        "• In-Memory RAG Vector Store grounds insights against statutory Ramsar, CPCB, and WHO standards."
    ], PURPLE_BIO)
]

for idx, (title, tag, pts, col) in enumerate(solutions_grid):
    col_idx = idx % 2
    row_idx = idx // 2
    c_left = Inches(0.8 + col_idx * 5.95)
    c_top  = Inches(1.8 + row_idx * 2.5)
    
    add_card(slide3, c_left, c_top, Inches(5.75), Inches(2.35), title=title, category_tag=tag, bg_color=OCEAN_CARD, border_color=col, border_width=1.5)
    add_bullet_list(slide3, c_left + Inches(0.2), c_top + Inches(0.8), Inches(5.35), Inches(1.45), pts, font_size=9.0, spacing=Pt(3))

add_footer_telemetry(slide3)

# ==============================================================================
# SLIDE 4: ORBITAL PIPELINE NODES (WITH EMBEDDED IMAGES AT COLUMN BOTTOMS)
# ==============================================================================
slide4 = create_slide(prs)
add_slide_header(slide4, 4, "Execution Telemetry", "Live Multi-Sensor Pipeline & Data Flow", "How raw optical, microwave radar, and topographic reflections transform into quantified volumetric & ecological intelligence.")

pipeline_nodes = [
    ("NODE 01", "Multi-Sensor STAC Ingest", [
        "INPUT: AOI BBox + Target Epochs",
        "• Query Microsoft Planetary STAC",
        "• Ingest S2 L2A BOA Optical",
        "• Ingest S1 RTC C-Band SAR Radar",
        "• Ingest Copernicus DEM GLO-30",
        "OUTPUT: Calibrated Multi-Sensor Grids"
    ], OCEAN_PANEL, AQUA_BRIGHT, IMG_NODE1, "MULTI-SENSOR STAC"),
    ("NODE 02", "Radiometric & Radar Math", [
        "INPUT: Spectral Bands + Radar Backscatter",
        "• NDWI = (B03-B08)/(B03+B08)",
        "• NDTI = (B04-B03)/(B04+B03) (Turbidity)",
        "• NDCI = (B05-B04)/(B05+B04) (Chl-a)",
        "• SAR Specular σ⁰_VV < -16 dB",
        "OUTPUT: 12-D Radiometric Tensor"
    ], OCEAN_PANEL, TEAL_BIO, IMG_NODE2, "RADIOMETRIC MATH"),
    ("NODE 03", "3D Hypsometry & Diff Mask", [
        "INPUT: DEM Contours + Water Extent",
        "• Hypsometric Volume: V = ∫ A(z)dz",
        "• Tri-state change matrix {-1, 0, +1}",
        "• Depth strata: 0-2m, 2-5m, >5m",
        "• Annual STAC time-series audit",
        "OUTPUT: km² / m³ Metrics & Diff Raster"
    ], OCEAN_PANEL, SOLAR_AMBER, IMG_NODE3, "3D HYPSOMETRY & DIFF"),
    ("NODE 04", "Dual AI Reasoning & RAG", [
        "INPUT: 12-D Tensor + Ground Photos",
        "• Sovereign Local RAG (Ramsar/CPCB)",
        "• Few-Shot Land-Cover Classification",
        "• Gemini 3.7: Deep CoT synthesis",
        "• Google Search & Drone Vision",
        "OUTPUT: Policy & Audit JSON Bundle"
    ], OCEAN_PANEL, CORAL_CRIMSON, IMG_NODE4, "SOVEREIGN RAG + AI")
]

for idx, (node, title, steps, bg_c, border_c, img_path, img_lbl) in enumerate(pipeline_nodes):
    left = Inches(0.8 + idx * 3.0)
    add_card(slide4, left, Inches(1.75), Inches(2.8), Inches(3.2), title=title, category_tag=node, bg_color=OCEAN_CARD, border_color=border_c, border_width=1.5)
    add_bullet_list(slide4, left + Inches(0.18), Inches(2.45), Inches(2.44), Inches(2.3), steps, font_size=8.5, spacing=Pt(2))
    add_image_with_frame(slide4, img_path, left, Inches(5.05), Inches(2.8), Inches(1.85), label=img_lbl, border_color=border_c)

add_footer_telemetry(slide4)

# ==============================================================================
# SLIDE 5: HIGH-COMPLEXITY TECHNICAL ARCHITECTURE DIAGRAM
# ==============================================================================
slide5 = create_slide(prs)
add_slide_header(slide5, 5, "Technical Approach", "End-to-End System Architecture & Technical Blueprint", "Comprehensive architecture connecting STAC multi-sensor nodes, 12-D tensor engine, few-shot classifier, local RAG vector store, and Gemini intelligence suite.")

add_image_with_frame(slide5, IMG_ARCH, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.15), label="SYSTEM ARCHITECTURE // COMPLETE TECHNICAL BLUEPRINT", border_color=OCEAN_BORDER_HI)

add_footer_telemetry(slide5)

# ==============================================================================
# SLIDE 6: REMOTE SENSING PHYSICS & BIO-OPTICAL MODELS
# ==============================================================================
slide6 = create_slide(prs)
add_slide_header(slide6, 6, "Remote Sensing Physics", "Multi-Sensor Electromagnetics & Bio-Optical Formulations", "Physical wavelength properties, radar backscatter, 3D hypsometry, and bio-optical water quality models.")

# Left Table Card: Multi-Sensor Bands + Math
add_card(slide6, Inches(0.8), Inches(1.8), Inches(6.4), Inches(4.9), title="Multi-Sensor Constellation Physics", category_tag="ELECTROMAGNETIC SPECTRUM & RADAR", border_color=AQUA_BRIGHT)

band_headers = ["Sensor / Band", "Wavelength / Freq", "Resolution", "Physical Role in AquaSense"]
band_rows = [
    ["S2 B03 (Green)", "560 nm", "10 m", "Primary NDWI numerator (high water reflection)"],
    ["S2 B04 (Red)", "665 nm", "10 m", "Chlorophyll absorption & NDTI turbidity numerator"],
    ["S2 B05 (Red-Edge)", "705 nm", "20 m", "Phytoplankton bloom scattering peak (NDCI)"],
    ["S2 B08 (NIR)", "842 nm", "10 m", "Pure water absorption (NDWI denominator)"],
    ["S1 C-Band SAR", "5.405 GHz (5.6cm)", "10 m", "All-weather cloud penetration via specular backscatter"],
    ["Copernicus DEM", "GLO-30 Topography", "30 m", "Hypsometric depth integration: V = ∫ A(z)dz"]
]
add_table_grid(slide6, Inches(1.0), Inches(2.6), Inches(6.0), Inches(2.7), band_headers, band_rows, [Inches(1.4), Inches(1.1), Inches(0.9), Inches(2.6)])

# Lower Mathematical Formula Box
add_rect(slide6, Inches(1.0), Inches(5.45), Inches(6.0), Inches(1.1), OCEAN_DEEP, OCEAN_BORDER)
tx_f = slide6.shapes.add_textbox(Inches(1.15), Inches(5.5), Inches(5.7), Inches(1.0))
tf_f = tx_f.text_frame
tf_f.word_wrap = True
tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
pf1 = tf_f.paragraphs[0]
pf1.text = "NDTI = (B04 - B03)/(B04 + B03) [Turbidity NTU]  |  NDCI = (B05 - B04)/(B05 + B04) [Chl-a µg/L]"
pf1.font.size = Pt(8.0)
pf1.font.bold = True
pf1.font.color.rgb = AQUA_BRIGHT
pf1.font.name = FONT_MONO

pf2 = tf_f.add_paragraph()
pf2.text = "3D Hypsometry: V(h) = sum 1/3(A_i + sqrt(A_i * A_{i+1}) + A_{i+1}) * Delta_h  |  WQI in [0, 100]"
pf2.font.size = Pt(8.0)
pf2.font.color.rgb = TEXT_FOAM
pf2.font.name = FONT_MONO
pf2.space_before = Pt(2)

# Right: High-Resolution Embedded NDWI Split Raster Image
add_image_with_frame(slide6, IMG_SPECTRAL, Inches(7.5), Inches(1.8), Inches(5.0), Inches(3.7), label="TRUE-COLOR VS NDWI + SAR WATER MASK")

# Right Lower: Image Caption Card
add_card(slide6, Inches(7.5), Inches(5.6), Inches(5.0), Inches(1.1), title="Visualizing Hydrological Inundation", category_tag="MULTI-SENSOR FUSED MASKS", bg_color=OCEAN_CARD, border_color=TEAL_BIO)
add_bullet_list(slide6, Inches(7.7), Inches(6.05), Inches(4.6), Inches(0.6), [
    "• Left: Sentinel-2 True Color (TCI) RGB optical view.",
    "• Right: High-contrast NDWI + SAR radar fused water mask with 3D depth contours."
], font_size=8.5, spacing=Pt(1))

add_footer_telemetry(slide6)

# ==============================================================================
# SLIDE 7: CARTOGRAPHIC OBSERVATORY UI BREAKDOWN
# ==============================================================================
slide7 = create_slide(prs)
add_slide_header(slide7, 7, "Product Interface", "AquaSense Cartographic Observatory & Streamlined Layout", "High-density scientific user interface delivering synchronized spatial, spectral, volumetric, and temporal insights.")

# UI Mockup Layout Simulation (3 Panes)
add_card(slide7, Inches(0.8), Inches(1.8), Inches(3.2), Inches(4.9), title="Telemetry & Hydro-Depth", category_tag="LEFT: CONFIG & VOLUMETRICS", bg_color=OCEAN_CARD, border_color=OCEAN_BORDER)
add_bullet_list(slide7, Inches(1.0), Inches(2.6), Inches(2.8), Inches(3.9), [
    "• Hero Setup & Collapsible Rail:",
    "  Prominent initialization console animating into dock/undock sidebar.",
    "",
    "• Interactive 8-Point BBOX Map:",
    "  Draggable anchor handles for custom spatial bounds [minLon, minLat, maxLon, maxLat].",
    "",
    "• 3D Hydro-Depth & Bio-Optics:",
    "  Recharts Area-Elevation volume curves, NTU turbidity, Chl-a, and CDOM.",
    "",
    "• Spectral Cutoff Tuning Sliders:",
    "  Optical NDWI cutoff & SAR dB threshold."
], font_size=9.0, spacing=Pt(3))

add_card(slide7, Inches(4.2), Inches(1.8), Inches(4.9), Inches(4.9), title="Observatory Canvas Stage", category_tag="CENTER: STREAMLINED 5-CORE MODES", bg_color=OCEAN_CARD_HI, border_color=AQUA_BRIGHT, border_width=1.5)
add_bullet_list(slide7, Inches(4.4), Inches(2.6), Inches(4.5), Inches(3.9), [
    "• 5 Streamlined Core View Modes:",
    "  1. [TRUE COLOR SWIPE] - Baseline RGB swipe",
    "  2. [NDWI WATER INDEX] - Pure water delineation",
    "  3. [SAR RADAR C-BAND] - All-weather microwave",
    "  4. [ALL-WEATHER FUSION] - Dual-sensor cross-val",
    "  5. [CHANGE DELTA] - Tri-state loss/gain matrix",
    "",
    "• Bio-Optics & 3D Dropdown Selector:",
    "  3D Depth ($m^3$), Turbidity, Chl-a, and CDOM.",
    "",
    "• Interactive Mode Physics Guide HUD:",
    "  Displays active formula, mini-summary, physical mechanism, and diagnostic utility.",
    "",
    "• 7 Scientific LUT Palettes with Needle Bar"
], font_size=9.0, spacing=Pt(2.5))

add_card(slide7, Inches(9.3), Inches(1.8), Inches(3.2), Inches(4.9), title="Quantification & Dual AI", category_tag="RIGHT: METRICS & REASONING", bg_color=OCEAN_CARD, border_color=OCEAN_BORDER)
add_bullet_list(slide7, Inches(9.5), Inches(2.6), Inches(2.8), Inches(3.9), [
    "• Hydrological Change Matrix:",
    "  Baseline vs. Target area (km²) & volume (MCM).",
    "  Net Change: -0.64 km² (-15.53%).",
    "",
    "• Multi-Year Recharts Trajectory:",
    "  Annual intermediate STAC samples.",
    "",
    "• Dual AI Intelligence Tabs:",
    "  1. Gemini 3.7 CoT Deep Synthesis",
    "  2. Search & Maps Grounding",
    "  3. Drone Multimodal Vision",
    "  4. Sovereign 12-D Tensor & RAG"
], font_size=9.0, spacing=Pt(3))

add_footer_telemetry(slide7)

# ==============================================================================
# SLIDE 8: DUAL AI ARCHITECTURE: CLOUD GEMINI 3.7 & SOVEREIGN 12-D RAG
# ==============================================================================
slide8 = create_slide(prs)
add_slide_header(slide8, 8, "Dual AI Architecture", "Cloud Gemini 3.7 & Sovereign 12-D Local RAG Engine", "Combining cloud multimodal reasoning with sovereign in-memory vector RAG and few-shot classification.")

# Left side: Cloud Gemini Modes
left_modes = [
    ("🧠 DEEP ECOLOGICAL REASONING", "gemini-3.7-flash (CoT Thinking)", [
        "• Formulates structured hydrological trajectory reports.",
        "• Isolates primary drivers (urban encroachment, rainfall deficits).",
        "• Assesses flood attenuation & groundwater ecosystem services."
    ], AQUA_BRIGHT),
    ("🔍 SEARCH & MAPS GROUNDING", "gemini-3.5-flash + Tools", [
        "• Integrates live Google Search for municipal restoration orders.",
        "• Google Maps spatial database identifies surrounding IT corridors.",
        "• Returns verifiable, clickable grounding citations."
    ], TEAL_BIO),
    ("🧬 SOVEREIGN 12-D LOCAL RAG", "In-Memory Vector Store (<35ms)", [
        "• Ingests Ramsar criteria, CPCB Class A-E standards, WHO limits.",
        "• Cosine Top-K retrieval & deterministic assessment synthesis.",
        "• 100% sovereign, private, and zero cloud API dependency."
    ], PURPLE_BIO)
]

for idx, (title, model, pts, col) in enumerate(left_modes):
    top = Inches(1.8 + idx * 1.68)
    add_card(slide8, Inches(0.8), top, Inches(6.4), Inches(1.55), title=title, category_tag=model, border_color=col)
    add_bullet_list(slide8, Inches(1.0), top + Inches(0.75), Inches(6.0), Inches(0.75), pts, font_size=8.5, spacing=Pt(1))

# Right side: Embedded Drone Vision Image + Card
add_image_with_frame(slide8, IMG_DRONE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(3.4), label="MULTIMODAL FIELD VALIDATION")

add_card(slide8, Inches(7.5), Inches(5.3), Inches(5.0), Inches(1.4), title="📸 Drone & Ground-Truth Photo Understanding", category_tag="GEMINI 3.7 FLASH (MULTIMODAL VISION)", bg_color=OCEAN_CARD, border_color=CORAL_CRIMSON)
add_bullet_list(slide8, Inches(7.7), Inches(5.95), Inches(4.6), Inches(0.7), [
    "• Ingests field photos, drone captures, or handheld phone images.",
    "• Evaluates water turbidity, surface algae, and hyacinth mats.",
    "• Produces estimated Few-Shot class confidence breakdowns."
], font_size=8.5, spacing=Pt(1))

add_footer_telemetry(slide8)

# ==============================================================================
# SLIDE 9: 12-DIMENSIONAL TENSOR PHYSICS & FEW-SHOT ML
# ==============================================================================
slide9 = create_slide(prs)
add_slide_header(slide9, 9, "Sovereign Machine Learning", "12-Dimensional Radiometric Extractor & Few-Shot ML", "Deterministic mathematical feature tensors feeding localized classifiers and statutory RAG retrieval.")

# Left Card: 12-D Feature Tensor
add_card(slide9, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.9), title="12-D Radiometric Tensor Architecture", category_tag="SPECTRAL12D EXTRACTOR ENGINE", border_color=AQUA_BRIGHT)

foundation_info = [
    "1. 12-DIMENSIONAL FEATURE TENSOR DEFINITION:",
    "   e = [B02, B03, B04, B08, NDWI, MNDWI, NDTI, NDCI, VV, VH, Delta_DEM, WQI]",
    "   • Combines visible bands, optical water indices, SAR radar backscatter, topographic elevation, and water quality index.",
    "",
    "2. DETERMINISTIC RADIOMETRIC PHYSICS:",
    "   • B02-B08: Calibrated surface reflectance (0-1).",
    "   • NDWI & MNDWI: Normalized Difference Water Indices.",
    "   • NDTI: Normalized Difference Turbidity Index (B04-B03)/(B04+B03).",
    "   • NDCI: Normalized Difference Chlorophyll Index (B05-B04)/(B05+B04).",
    "   • VV & VH: Sentinel-1 C-band backscatter in decibels (dB).",
    "   • Delta_DEM: Copernicus GLO-30 topographic slope buffer.",
    "   • WQI: Aggregated bio-optical ecological score (0-100).",
    "",
    "3. ZERO-DEPENDENCY EDGE INFERENCE:",
    "   • Extracted in pure TypeScript/Node.js without external GPU."
]
add_bullet_list(slide9, Inches(1.0), Inches(2.6), Inches(5.3), Inches(3.9), foundation_info, font_size=9.0, spacing=Pt(2.5))

# Right Card: FewShotClassifier Mathematics
add_card(slide9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9), title="FewShotClassifier & RAG Retrieval", category_tag="RAPID ADAPTATION OVER 12-D TENSORS", border_color=TEAL_BIO)

fewshot_math = [
    "1. L2-REGULARIZED LOGISTIC REGRESSION:",
    "   Loss = -1/N sum [ w+ y log(sigma(z)) + w- (1-y) log(1-sigma(z)) ] + 1/(2C) ||w||²",
    "   • Inverse Class Frequency Weighting: w+ = N/(2*N+), w- = N/(2*N-).",
    "   • Output Classes: water, wetland, built_up.",
    "",
    "2. COSINE-WEIGHTED k-NEAREST NEIGHBORS (k-NN):",
    "   Distance = 1.0 - (x . r_j) / (||x|| * ||r_j||)",
    "   • Voting Weights: w_j = 1.0 / (dist_j + 1e-5).",
    "",
    "3. IN-MEMORY RAG STATUTORY RETRIEVAL:",
    "   • TF-IDF Cosine semantic similarity over Ramsar criteria, CPCB Class A-E benchmarks, and WHO limits.",
    "   • Top-K retrieved chunks feed structured markdown synthesizer."
]
add_bullet_list(slide9, Inches(7.0), Inches(2.6), Inches(5.3), Inches(3.9), fewshot_math, font_size=9.0, spacing=Pt(2.5))

add_footer_telemetry(slide9)

# ==============================================================================
# SLIDE 10: QUANTITATIVE CASE STUDY: PALLIKARANAI MARSHLAND
# ==============================================================================
slide10 = create_slide(prs)
add_slide_header(slide10, 10, "Field Validation", "Pallikaranai Ramsar Marshland Multi-Sensor Audit", "Quantifying surface water shrinkage, volumetric retention loss, and eutrophication across 2019-2025.")

# Top 4 Quantitative Metric Tiles
case_stats = [
    ("4.12 km²", "2019 BASELINE WATER (T0)", "Post-monsoon water surface extent", AQUA_BRIGHT),
    ("3.48 km²", "2025 TARGET WATER (T1)", "Recent Sentinel-2/1 pass extent", AQUA_DEEP),
    ("-0.64 km²", "ABSOLUTE NET CHANGE", "Permanent water surface lost", CORAL_CRIMSON),
    ("-15.53%", "RELATIVE WETLAND LOSS", "Severe hydrological constriction", SOLAR_AMBER)
]

for idx, (val, title, desc, col) in enumerate(case_stats):
    left = Inches(0.8 + idx * 3.0)
    card = add_rect(slide10, left, Inches(1.8), Inches(2.8), Inches(1.2), OCEAN_CARD, col, border_width=1.5)
    
    tx = slide10.shapes.add_textbox(left + Inches(0.15), Inches(1.85), Inches(2.5), Inches(1.1))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = val
    p0.font.size = Pt(20)
    p0.font.bold = True
    p0.font.color.rgb = col
    p0.font.name = FONT_TITLE
    
    p1 = tf.add_paragraph()
    p1.text = title
    p1.font.size = Pt(8)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_GLACIAL
    p1.font.name = FONT_MONO
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = TEXT_SEABED

# Left Column: Spatial Findings & Ground-Truth
add_card(slide10, Inches(0.8), Inches(3.2), Inches(5.7), Inches(3.5), title="Spatial Findings & Bio-Optical Audit", category_tag="MULTI-SENSOR GROUND-TRUTH CONFIRMATION", border_color=CORAL_CRIMSON)

findings = [
    "• Northern Perimeter Shrinkage: The Diff Mask reveals intense signal red loss along the Velachery-Tambaram IT corridor expansion.",
    "• Volumetric Storage Loss: 3D DEM Bathymetry reveals a drop from 12.8 MCM (2019) to 10.4 MCM (2025) (-18.75% storage volume).",
    "• Bio-Optical Water Quality: Turbidity rose from 14.2 NTU to 28.6 NTU; Chlorophyll-a rose to 19.4 µg/L (Eutrophic Alert, WQI: 68).",
    "• Multi-Sensor Verification: Baseline S2A_MSIL2A (0% cloud) vs. Target S2B + S1 C-SAR (all-weather radar confirmed)."
]
add_bullet_list(slide10, Inches(1.0), Inches(3.9), Inches(5.3), Inches(2.6), findings, font_size=9.0, spacing=Pt(3))

# Right Column: Annual Longitudinal Trend Data Table
add_card(slide10, Inches(6.8), Inches(3.2), Inches(5.7), Inches(3.5), title="Longitudinal Time-Series & Volumetric Data", category_tag="INTERMEDIATE ANNUAL MULTI-SENSOR SAMPLING", border_color=TEAL_BIO)

trend_headers = ["Year", "Water (km²)", "Volume (MCM)", "WQI (0-100)", "Hydrological State"]
trend_rows = [
    ["2019", "4.12 km²", "12.8 MCM", "84 (Good)", "Baseline Post-Drought Recovery"],
    ["2020", "4.05 km²", "12.4 MCM", "82 (Good)", "Stable Pre-Monsoon Capacity"],
    ["2021", "4.38 km²", "14.1 MCM", "88 (Good)", "Extreme Monsoon Inundation Peak"],
    ["2022", "3.89 km²", "11.9 MCM", "76 (Fair)", "Construction Encroachment Phase"],
    ["2023", "3.62 km²", "11.0 MCM", "71 (Fair)", "Summer Dry-Down & Constriction"],
    ["2024", "3.51 km²", "10.6 MCM", "69 (Moderate)", "Fragmented Southern Retention"],
    ["2025", "3.48 km²", "10.4 MCM", "68 (Moderate)", "Current Multi-Sensor Audit (-15.5%)"]
]
add_table_grid(slide10, Inches(7.0), Inches(3.9), Inches(5.3), Inches(2.6), trend_headers, trend_rows, [Inches(0.8), Inches(1.1), Inches(1.1), Inches(1.0), Inches(1.3)])

add_footer_telemetry(slide10)

# ==============================================================================
# SLIDE 11: GLOBAL IMPACT & STRATEGIC SCALING ROADMAP
# ==============================================================================
slide11 = create_slide(prs)
add_slide_header(slide11, 11, "Strategic Impact & Roadmap", "Global Stakeholder Applications & Strategic Scaling Roadmap", "Transforming satellite observation into verifiable conservation policy, smart city resilience, and planetary hydro-intelligence.")

# LEFT SECTION: 4 Stakeholder Applications (2x2 Grid)
add_card(slide11, Inches(0.8), Inches(1.75), Inches(5.7), Inches(5.0), title="Target Stakeholders & Real-World Impact", category_tag="GLOBAL APPLICATION DOMAINS", bg_color=OCEAN_CARD, border_color=OCEAN_BORDER_HI)

stakeholders_mini = [
    ("🏛️ GOVERNANCE & POLICY", "WETLAND AUTHORITIES", [
        "• Ramsar & Biodiversity compliance reporting.",
        "• Automated alerts on illegal landfilling.",
        "• Court-admissible proof (hash: 0x8a92f02c)."
    ], AQUA_BRIGHT),
    ("🌱 CONSERVATION NGOS", "ECOLOGISTS & ADVOCACY", [
        "• Continuous remote monitoring of basins.",
        "• Field photo ground-truthing portal.",
        "• Multi-year trend data for climate grants."
    ], TEAL_BIO),
    ("🏙️ SMART CITIES & PLANNERS", "FLOOD ATTENUATION", [
        "• 3D volumetric water holding capacity checks.",
        "• Critical urban flood buffer preservation.",
        "• City emergency dashboard integration."
    ], SOLAR_AMBER),
    ("💼 ESG & BLUE CARBON", "SUSTAINABILITY AUDIT", [
        "• Independent satellite MRV verification.",
        "• Auditable wetland carbon credit validation.",
        "• Corporate water stewardship reporting."
    ], CORAL_CRIMSON)
]

for idx, (title, tag, pts, col) in enumerate(stakeholders_mini):
    col_idx = idx % 2
    row_idx = idx // 2
    c_left = Inches(1.0 + col_idx * 2.7)
    c_top  = Inches(2.45 + row_idx * 2.05)
    
    add_rect(slide11, c_left, c_top, Inches(2.55), Inches(1.95), OCEAN_DEEP, col, border_width=1)
    
    txBox = slide11.shapes.add_textbox(c_left + Inches(0.1), c_top + Inches(0.1), Inches(2.35), Inches(1.75))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(8.5)
    p0.font.bold = True
    p0.font.color.rgb = col
    p0.font.name = FONT_TITLE
    
    p1 = tf.add_paragraph()
    p1.text = tag
    p1.font.size = Pt(7.0)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_SEABED
    p1.font.name = FONT_MONO
    p1.space_after = Pt(3)
    
    for pt in pts:
        p_pt = tf.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(8.0)
        p_pt.font.color.rgb = TEXT_FOAM
        p_pt.space_after = Pt(2)

# RIGHT SECTION: 4 Strategic Scaling Roadmap Phases (Vertical Stack)
add_card(slide11, Inches(6.8), Inches(1.75), Inches(5.7), Inches(5.0), title="Strategic Scaling Roadmap & Milestones", category_tag="PLANETARY DEPLOYMENT PHASES", bg_color=OCEAN_CARD, border_color=TEAL_BIO)

roadmap_phases = [
    ("PHASE 1: CURRENT PRODUCTION", "BENCHMARK DEPLOYMENT", "✔ Multi-sensor STAC, 5-core view modes, 3D bathymetry (m³), bio-optics (NDTI/NDCI), 12-D RAG engine, and Pallikaranai audit.", TEAL_BIO),
    ("PHASE 2: Q4 2025", "MULTI-BASIN EXPANSION", "• Global STAC webhook triggers, automated GeoJSON vector boundary export, and Landsat thermal infrared integration.", AQUA_BRIGHT),
    ("PHASE 3: Q1 2026", "ACTIVE LEARNING ML", "• Few-shot patch labeling UI, edge WebAssembly deployment, and citizen science mobile photo ground-truthing portal.", SOLAR_AMBER),
    ("PHASE 4: 2027 VISION", "PLANETARY HYDRO-NETWORK", "• Pan-India & Global Wetland Atlas, real-time flood surge forecasting, open public API for NGOs, automated Blue Carbon registry.", CORAL_CRIMSON)
]

for idx, (phase, tag, desc, col) in enumerate(roadmap_phases):
    p_top = Inches(2.45 + idx * 1.02)
    add_rect(slide11, Inches(7.0), p_top, Inches(5.3), Inches(0.92), OCEAN_DEEP, col, border_width=1)
    
    txBox = slide11.shapes.add_textbox(Inches(7.15), p_top + Inches(0.08), Inches(5.0), Inches(0.78))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = f"{phase}  //  {tag}"
    p0.font.size = Pt(8.5)
    p0.font.bold = True
    p0.font.color.rgb = col
    p0.font.name = FONT_MONO
    
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(8.0)
    p1.font.color.rgb = TEXT_FOAM
    p1.space_before = Pt(2)

add_footer_telemetry(slide11)

# ==============================================================================
# SLIDE 12: LIVE OBSERVATORY PRODUCT OUTPUT (WITH EMBEDDED DASHBOARD SCREENSHOT)
# ==============================================================================
slide12 = create_slide(prs)
add_slide_header(slide12, 12, "Live Product Output", "AquaSense Live Planetary Observatory: Verified Operational Output", "Direct screen capture of the operational multi-sensor dashboard running on Sentinel-2 MSI, Sentinel-1 C-SAR, and DEM hypsometry.")

# Large High-Resolution Embedded Application Screenshot
add_image_with_frame(slide12, IMG_SCREENSHOT, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2), label="AQUASENSE OPERATIONAL MULTI-SENSOR DASHBOARD // LIVE RUN OUTPUT", border_color=BORDER_AQUA)

add_footer_telemetry(slide12, text="AQUASENSE PLANETARY OBSERVATORY • LIVE RUN OUTPUT • ALL-WEATHER RADAR FUSION • 3D BATHYMETRY • 12-D RAG ENGINE")

# ==============================================================================
# SAVE PRESENTATION TO LOCAL WORKSPACE DIRECTORY
# ==============================================================================
output_filename = "AquaSense_Pitch_Deck.pptx"
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), output_filename)
try:
    prs.save(output_path)
    print(f"[SUCCESS] AquaSense presentation generated successfully: {output_path}")
except PermissionError:
    alt_filename = "AquaSense_Pitch_Deck_v2.pptx"
    alt_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), alt_filename)
    prs.save(alt_path)
    print(f"[SUCCESS] AquaSense presentation saved to: {alt_path}")
